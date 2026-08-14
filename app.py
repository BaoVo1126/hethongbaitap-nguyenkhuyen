import os
import re
import json
import sqlite3
import random
import hashlib
import uuid
import subprocess
from datetime import datetime, date
import requests 
from flask import Flask, g, render_template, request, redirect, url_for, flash, abort, send_from_directory
from werkzeug.utils import secure_filename

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

os.environ["FLAGS_use_mkldnn"] = "0"
os.environ["FLAGS_enable_pir_api"] = "0"

try:
    import pymupdf as fitz  
    PDF_AVAILABLE = True
except ImportError:
    try:
        import fitz
        PDF_AVAILABLE = True
    except ImportError:
        PDF_AVAILABLE = False

try:
    from paddleocr import PaddleOCR
    PADDLE_AVAILABLE = True
    _ocr = PaddleOCR(use_textline_orientation=True, lang='en')
except ImportError:
    PADDLE_AVAILABLE = False
    _ocr = None

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False


OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434/api/generate")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "llama3.2:latest")  

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

if os.environ.get("RENDER") or os.environ.get("VERCEL"):
    DATA_DIR = "/tmp"
else:
    DATA_DIR = BASE_DIR

DB_PATH = os.path.join(DATA_DIR, "center.db")
UPLOAD_DIR = os.path.join(DATA_DIR, "uploads")
DOCS_DIR = os.path.join(UPLOAD_DIR, "documents")
CONVERTED_DIR = os.path.join(UPLOAD_DIR, "_converted")
EXTRACTED_IMAGES_DIR = os.path.join(DATA_DIR, "static", "extracted_images")

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(DOCS_DIR, exist_ok=True)
os.makedirs(CONVERTED_DIR, exist_ok=True)
os.makedirs(EXTRACTED_IMAGES_DIR, exist_ok=True)

CENTER_NAME = "Trung Tâm Nguyễn Khuyến"
CENTER_INFO = {
    "tagline": "Hệ thống luyện tập & kiểm tra trực tuyến",
    "subjects": ["Toán", "Hoá", "Lý", "Anh văn"],
    "description": (
        "Trung tâm tổ chức các buổi luyện tập và kiểm tra theo từng chuyên đề. "
        "Học sinh làm bài trực tiếp trên hệ thống, được chấm điểm và xem lại lỗi sai ngay sau khi nộp."
    ),
}

TEACHER_PASSWORD = os.environ.get("TEACHER_PASSWORD", "1111")

app = Flask(__name__)
app.config["SECRET_KEY"] = "demo-secret-key-doi-khi-deploy-that"
app.config["MAX_CONTENT_LENGTH"] = 32 * 1024 * 1024  # 32MB

RANDOMIZE_EXAM = True


@app.template_filter('from_json')
def from_json_filter(value):
    if not value:
        return []
    try:
        return json.loads(value)
    except Exception:
        return []


def get_db():
    if "db" not in g:
        g.db = sqlite3.connect(DB_PATH)
        g.db.row_factory = sqlite3.Row
        g.db.execute("PRAGMA foreign_keys = ON")
    return g.db


@app.teardown_appcontext
def close_db(exception=None):
    db = g.pop("db", None)
    if db is not None:
        db.close()


def init_db():
    db = sqlite3.connect(DB_PATH)
    db.executescript(
        """
        CREATE TABLE IF NOT EXISTS exams (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            subject TEXT,
            source_image TEXT,
            raw_ocr_text TEXT,
            status TEXT NOT NULL DEFAULT 'draft',
            available_date TEXT,
            time_limit_minutes INTEGER,
            created_at TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS questions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            exam_id INTEGER NOT NULL REFERENCES exams(id) ON DELETE CASCADE,
            order_index INTEGER NOT NULL,
            question_text TEXT NOT NULL,
            question_type TEXT NOT NULL DEFAULT 'mcq',
            option_a TEXT NOT NULL,
            option_b TEXT NOT NULL,
            option_c TEXT NOT NULL,
            option_d TEXT NOT NULL,
            correct_answer TEXT,
            correct_answer_text TEXT,
            image_urls TEXT DEFAULT '[]'
        );

        CREATE TABLE IF NOT EXISTS tests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            exam_id INTEGER NOT NULL REFERENCES exams(id) ON DELETE CASCADE,
            student_name TEXT NOT NULL,
            student_code TEXT NOT NULL,
            question_order TEXT NOT NULL,
            option_maps TEXT NOT NULL,
            created_at TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS submissions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            test_id INTEGER NOT NULL REFERENCES tests(id) ON DELETE CASCADE,
            answers TEXT NOT NULL,
            score INTEGER NOT NULL,
            total INTEGER NOT NULL,
            submitted_at TEXT NOT NULL
        );

        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            original_filename TEXT NOT NULL,
            stored_filename TEXT NOT NULL,
            subject TEXT,
            uploaded_at TEXT NOT NULL
        );
        """
    )
    cols = [r[1] for r in db.execute("PRAGMA table_info(exams)").fetchall()]
    if "available_date" not in cols:
        db.execute("ALTER TABLE exams ADD COLUMN available_date TEXT")
    if "time_limit_minutes" not in cols:
        db.execute("ALTER TABLE exams ADD COLUMN time_limit_minutes INTEGER")

    doc_cols = [r[1] for r in db.execute("PRAGMA table_info(documents)").fetchall()]
    if "subject" not in doc_cols:
        db.execute("ALTER TABLE documents ADD COLUMN subject TEXT")
        
    q_cols = [r[1] for r in db.execute("PRAGMA table_info(questions)").fetchall()]
    if "image_urls" not in q_cols:
        db.execute("ALTER TABLE questions ADD COLUMN image_urls TEXT DEFAULT '[]'")
    if "question_type" not in q_cols:
        db.execute("ALTER TABLE questions ADD COLUMN question_type TEXT NOT NULL DEFAULT 'mcq'")
    if "correct_answer_text" not in q_cols:
        db.execute("ALTER TABLE questions ADD COLUMN correct_answer_text TEXT")

    db.commit()
    db.close()


# --------------------------------------------------------------------------
# Pipeline: Deterministic document & Enhanced Image/Question Extractor
# --------------------------------------------------------------------------

def _clean_extracted_line(line: str) -> str:
    line = line.replace("\u00a0", " ")
    return re.sub(r"\s+", " ", line).strip()


def _join_wrapped_lines(text: str) -> str:
    lines = [_clean_extracted_line(x) for x in text.splitlines()]
    lines = [x for x in lines if x]
    return " ".join(lines).strip()


def _read_txt(file_path: str) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1258", "cp1252"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    return ""


def _read_docx(file_path: str, exam_id: int = None) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""

    doc = Document(file_path)
    parts = []
    
    if exam_id:
        doc_img_dir = os.path.join(EXTRACTED_IMAGES_DIR, str(exam_id))
        os.makedirs(doc_img_dir, exist_ok=True)
        img_index = 1
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    img_part = rel.target_part
                    img_ext = img_part.content_type.split('/')[-1]
                    if img_ext not in ['png', 'jpeg', 'jpg', 'gif', 'webp']:
                        img_ext = 'png'
                    img_name = f"docx_img_{img_index}_{uuid.uuid4().hex[:6]}.{img_ext}"
                    img_save_path = os.path.join(doc_img_dir, img_name)
                    with open(img_save_path, "wb") as f:
                        f.write(img_part.blob)
                    img_index += 1
                except Exception as e:
                    print(f"[DOCX Image Extract Warning] {e}")

    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _read_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""

    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _read_xlsx(file_path: str) -> str:
    if not OPENPYXL_AVAILABLE:
        return ""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    except Exception:
        return ""
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


LIBREOFFICE_BIN = None
for _candidate in ("libreoffice", "soffice"):
    from shutil import which as _which
    if _which(_candidate):
        LIBREOFFICE_BIN = _candidate
        break


def _convert_office_to_pdf(file_path: str):
    if not LIBREOFFICE_BIN:
        return None
    try:
        result = subprocess.run(
            [LIBREOFFICE_BIN, "--headless", "--norestore", "--convert-to", "pdf",
             "--outdir", CONVERTED_DIR, file_path],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            print(f"[LibreOffice Warning] {result.stderr}")
            return None
    except Exception as e:
        print(f"[LibreOffice Warning] {e}")
        return None

    base = os.path.splitext(os.path.basename(file_path))[0]
    candidate = os.path.join(CONVERTED_DIR, base + ".pdf")
    return candidate if os.path.exists(candidate) else None


def _looks_like_broken_font_rendering(text: str) -> bool:
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    if len(lines) < 8:
        return False
    lone_marks = sum(1 for ln in lines if len(ln) == 1 and not ln.isascii())
    return lone_marks >= 5 and (lone_marks / len(lines)) > 0.08


def _ocr_image_paddle(image) -> str:
    if not PADDLE_AVAILABLE or _ocr is None:
        return ""
    try:
        import numpy as np
        img_np = np.array(image)
        result = _ocr.predict(img_np)

        chunks = []
        for page_result in result or []:
            data = None
            if hasattr(page_result, "json"):
                try:
                    data = page_result.json
                    if callable(data):
                        data = data()
                except Exception:
                    data = None
            if data is None and hasattr(page_result, "res"):
                data = page_result.res
            if isinstance(data, str):
                try:
                    data = json.loads(data)
                except Exception:
                    pass
            if isinstance(data, dict):
                data = data.get("res", data)
                texts = data.get("rec_texts") or data.get("texts") or []
                chunks.extend(str(x) for x in texts if str(x).strip())
            elif isinstance(page_result, dict):
                texts = page_result.get("rec_texts") or page_result.get("texts") or []
                chunks.extend(str(x) for x in texts if str(x).strip())

        if chunks:
            return "\n".join(chunks)
    except Exception as e:
        print(f"[PaddleOCR Warning] {e}")

    return ""


def _ocr_image_tesseract(image) -> str:
    if not OCR_AVAILABLE:
        return ""
    try:
        return pytesseract.image_to_string(image, lang="vie+eng")
    except Exception:
        try:
            return pytesseract.image_to_string(image)
        except Exception:
            return ""


def _pdf_has_real_text(text: str) -> bool:
    if not text or len(text.strip()) < 40:
        return False
    question_hits = len(re.findall(r"(?m)^\s*(?:Câu\s+\d+|\d+\.\d+|\d+)\.\s+", text, re.IGNORECASE))
    word_hits = len(re.findall(r"[A-Za-zÀ-ỹ]{3,}", text))
    return question_hits >= 1 or word_hits >= 20


def extract_text_from_file(file_path: str, exam_id: int = None) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".txt":
        return _read_txt(file_path)

    if ext == ".docx":
        return _read_docx(file_path, exam_id=exam_id)

    if ext == ".pptx":
        return _read_pptx(file_path)

    if ext in {".xlsx", ".xls"}:
        return _read_xlsx(file_path)

    if ext == ".pdf" and PDF_AVAILABLE:
        doc = fitz.open(file_path)
        native_pages = []
        for page in doc:
            native_pages.append(page.get_text("text", sort=True))
        native_text = "\n".join(native_pages)

        if _pdf_has_real_text(native_text):
            return native_text

        if PDF2IMAGE_AVAILABLE:
            try:
                images = convert_from_path(file_path, dpi=220)
                ocr_pages = []
                for img in images:
                    page_text = _ocr_image_paddle(img)
                    if not page_text:
                        page_text = _ocr_image_tesseract(img)
                    ocr_pages.append(page_text)
                return "\n".join(ocr_pages)
            except Exception as e:
                print(f"[PDF OCR Warning] {e}")
        return native_text

    if ext in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tif", ".tiff"}:
        if not OCR_AVAILABLE and not PADDLE_AVAILABLE:
            return ""
        try:
            from PIL import Image
            image = Image.open(file_path)
            text = _ocr_image_paddle(image)
            return text or _ocr_image_tesseract(image)
        except Exception as e:
            print(f"[Image OCR Warning] {e}")
            return ""

    return ""


def _smart_join_fragments(parts) -> str:
    """Ghép các dòng bị bẻ (word-wrap) của một đáp án lại thành 1 chuỗi.
    Nếu một dòng kết thúc bằng dấu gạch nối do bị ngắt giữa từ, nối liền
    không thêm khoảng trắng."""
    cleaned = [p.strip() for p in parts if p and p.strip()]
    if not cleaned:
        return ""
    result = cleaned[0]
    for part in cleaned[1:]:
        if result.endswith("-") and not result.endswith("--"):
            result += part
        else:
            result += " " + part
    return re.sub(r"\s+", " ", result).strip()


# Chỉ nhận CHỮ HOA A/B/C/D làm nhãn phương án trắc nghiệm. Chữ thường a)/b)/
# c)/d) (thường dùng để đánh dấu các ý đúng/sai trong 1 câu, hoặc các ý a),
# b) của câu tự luận nhiều phần) KHÔNG được coi là phương án trắc nghiệm — nếu
# không sẽ nhận nhầm câu "đúng/sai" hoặc câu tự luận nhiều ý thành trắc nghiệm.
_OPTION_START_RE = re.compile(r"^(?:[\$]\s*)?([ABCD])(?![A-Za-zÀ-ỹ])[.\)~:]?\s*(.*)$")

# Nhận diện đầu 1 câu hỏi: "Câu 1.", "Bài 2.", hoặc số trần "19.1." — dấu chấm
# phải đứng NGAY sau số (không có khoảng trắng ở giữa) để không nhận nhầm 1
# mảnh số bị tách rời (ví dụ mẫu số của phân số y=1/7 bị PDF tách thành block
# riêng "7.") thành đầu 1 câu hỏi mới.
_QUESTION_START_RE = re.compile(r"^\s*(?:C[âa]u\s+|Bài\s+)?(\d+(?:\.\d+)?)\.\s*(.*)$", re.IGNORECASE)


def _parse_option_lines(text: str):
    """Trích 4 đáp án A/B/C/D từ văn bản thô — hỗ trợ mọi kiểu trình bày:
    mỗi đáp án 1 dòng riêng, 2 đáp án chung 1 dòng, đáp án dài bị word-wrap
    qua nhiều dòng, chữ cái đứng 1 mình rồi nội dung xuống dòng sau."""
    options = []
    current_letter = None
    current_parts = []

    for raw_line in text.splitlines():
        line = _clean_extracted_line(raw_line)
        if not line:
            continue
        if _looks_like_header_or_footer(line):
            continue

        m = _OPTION_START_RE.match(line)
        if m:
            if current_letter is not None:
                options.append((current_letter, _smart_join_fragments(current_parts)))
            current_letter = m.group(1).upper()
            val = m.group(2).strip()
            current_parts = [val] if val else []
        elif current_letter is not None:
            current_parts.append(line)

    if current_letter is not None:
        options.append((current_letter, _smart_join_fragments(current_parts)))

    return options


def _looks_like_header_or_footer(text: str) -> bool:
    """Nhận diện các dòng tiêu đề/chân trang/phân mục MANG TÍNH CẤU TRÚC
    (số trang, gạch đầu dòng, "PHẦN I/II/III.", "I./II.", hotline, khẩu hiệu
    trong ngoặc kép...) — cố tình KHÔNG dựa vào từ khóa riêng của 1 trường/1
    môn cụ thể nào, để dùng chung được cho bất kỳ đề của bất kỳ môn/giáo viên
    nào mà không cần chỉnh lại mỗi lần đổi mẫu đề."""
    t = text.lower().strip()
    if not t:
        return True

    if re.match(r"^[■•»›]", t):
        return True
    if re.match(r"^phần\s+[ivx]+\b", t):
        return True
    if re.match(r"^(?:i|ii|iii|iv|v)\.\s", t):
        return True
    if re.fullmatch(r"[-–—»\s]*\d+[-–—»\s]*", t):
        return True
    if re.search(r"\bhotline\b|\bsđt\b|\bđt\b\s*[:.]?\s*\d{7,}", t):
        return True
    # Khẩu hiệu/slogan chân trang thường để trong ngoặc kép nguyên 1 dòng
    if re.match(r'^["“‘].*["”’]$', text.strip()):
        return True

    return False


def _render_and_crop_drawings(page, exam_id, page_idx, doc_img_dir, top_ratio=0.08, bottom_ratio=0.93):
    """Phát hiện đồ thị/bảng biến thiên được VẼ BẰNG NÉT VECTOR (không phải
    ảnh bitmap nhúng sẵn) rồi crop thành PNG. Bỏ qua nét trang trí nhỏ (bullet,
    ô vuông điền đáp án...) và mọi nét nằm trong vùng đầu/cuối trang (logo, số
    trang, chân trang) — lọc theo VỊ TRÍ trên trang nên áp dụng được cho bất kỳ
    mẫu đề nào, không phụ thuộc nội dung chữ cụ thể."""
    extracted_imgs = []
    try:
        page_h = page.rect.height
        paths = page.get_drawings()
        if not paths:
            return extracted_imgs

        clusters = []
        for path in paths:
            r = path["rect"]
            if r.width < 15 and r.height < 15:
                continue
            if r.y1 < page_h * top_ratio or r.y0 > page_h * bottom_ratio:
                continue
            if r.width > page.rect.width * 0.9:
                continue

            found = False
            for c in clusters:
                if (abs(r.y0 - c["y1"]) < 30 or abs(r.y1 - c["y0"]) < 30
                        or (r.y0 >= c["y0"] - 5 and r.y1 <= c["y1"] + 5)):
                    c["x0"] = min(c["x0"], r.x0)
                    c["y0"] = min(c["y0"], r.y0)
                    c["x1"] = max(c["x1"], r.x1)
                    c["y1"] = max(c["y1"], r.y1)
                    c["count"] += 1
                    found = True
                    break
            if not found:
                clusters.append({"x0": r.x0, "y0": r.y0, "x1": r.x1, "y1": r.y1, "count": 1})

        for idx, c in enumerate(clusters):
            w, h = c["x1"] - c["x0"], c["y1"] - c["y0"]
            if w < 50 or h < 40 or w > page.rect.width * 0.95 or h > page.rect.height * 0.8:
                continue
            if c["count"] < 5:
                continue  # cụm quá ít nét, khó là 1 hình vẽ thật

            crop_rect = fitz.Rect(
                max(0, c["x0"] - 10), max(0, c["y0"] - 10),
                min(page.rect.width, c["x1"] + 10), min(page.rect.height, c["y1"] + 10),
            )
            pix = page.get_pixmap(clip=crop_rect, dpi=200)
            image_filename = f"p{page_idx}_vector{idx + 1}_{uuid.uuid4().hex[:6]}.png"
            image_save_path = os.path.join(doc_img_dir, image_filename)
            pix.save(image_save_path)
            web_url = f"/static/extracted_images/{exam_id}/{image_filename}"
            extracted_imgs.append((c["y0"], c["y1"], web_url))
    except Exception as e:
        print(f"[Vector Crop Warning] {e}")

    return extracted_imgs


def _make_question_dict(number, qtext, question_type, options, page, image_urls):
    return {
        "question_text": f"Câu {number}. {qtext}",
        "question_type": question_type,
        "option_a": options.get("A", ""),
        "option_b": options.get("B", ""),
        "option_c": options.get("C", ""),
        "option_d": options.get("D", ""),
        "correct_answer": None,
        "correct_answer_text": None,
        "source_page": page,
        "image_urls": json.dumps(image_urls),
    }


def _extract_questions_from_pdf_layout(file_path: str, exam_id: int = None):
    if not PDF_AVAILABLE:
        return []

    doc = fitz.open(file_path)
    doc_img_dir = None
    if exam_id:
        doc_img_dir = os.path.join(EXTRACTED_IMAGES_DIR, str(exam_id))
        os.makedirs(doc_img_dir, exist_ok=True)

    # ---- Bước 1: đọc TOÀN BỘ block + ảnh của mọi trang trước (gắn số trang
    # vào từng block/ảnh), thay vì xử lý xong-trang-nào-bỏ-trang-đó — để 1 câu
    # hỏi bị tràn qua trang sau (VD: đề bài ở cuối trang này, bảng biến thiên/
    # đồ thị nằm ở đầu trang kế tiếp) vẫn được ghép đúng và gán đúng ảnh.
    all_kept_blocks = []   # (page_idx, text, x0, y0, y1)
    all_page_images = []   # (page_idx, y0, y1, web_url)
    margin_counter = {}

    for page_idx, page in enumerate(doc, start=1):
        page_h = page.rect.height
        # Trang 1 thường có logo/tiêu đề lớn nên vùng "đầu trang" cần rộng hơn;
        # các trang sau chỉ cần chừa chỗ cho tiêu đề/số trang lặp lại mỏng.
        top_ratio = 0.20 if page_idx == 1 else 0.08
        bottom_ratio = 0.93

        raw_blocks = page.get_text("blocks", sort=True)
        for b in raw_blocks:
            x0, y0, x1, y1, text = b[0], b[1], b[2], b[3], (b[4] or "")
            if not text.strip():
                continue
            # Lọc theo VỊ TRÍ (không theo nội dung) để loại logo/tiêu đề lặp/số
            # trang/chân trang cho MỌI mẫu đề, không cần biết trước nó viết gì
            if y1 < page_h * top_ratio or y0 > page_h * bottom_ratio:
                continue
            all_kept_blocks.append((page_idx, text, x0, y0, y1))
            bucket = round(x0 / 4) * 4
            margin_counter[bucket] = margin_counter.get(bucket, 0) + 1

        if exam_id and doc_img_dir:
            for img_index, img in enumerate(page.get_images(full=True)):
                try:
                    xref = img[0]
                    img_rects = page.get_image_rects(xref)
                    if not img_rects:
                        continue
                    r = img_rects[0]
                    if r.width < 25 or r.height < 25:
                        continue  # icon/hoạ tiết trang trí quá nhỏ
                    if r.y1 < page_h * top_ratio or r.y0 > page_h * bottom_ratio:
                        continue  # logo/hoạ tiết nằm trong vùng đầu/cuối trang

                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    image_filename = f"p{page_idx}_img{img_index + 1}_{uuid.uuid4().hex[:6]}.{image_ext}"
                    image_save_path = os.path.join(doc_img_dir, image_filename)
                    with open(image_save_path, "wb") as f:
                        f.write(image_bytes)
                    web_url = f"/static/extracted_images/{exam_id}/{image_filename}"
                    all_page_images.append((page_idx, r.y0, r.y1, web_url))
                except Exception as e:
                    print(f"[PDF Image Extraction Warning] {e}")

            for (iy0, iy1, url) in _render_and_crop_drawings(page, exam_id, page_idx, doc_img_dir, top_ratio, bottom_ratio):
                all_page_images.append((page_idx, iy0, iy1, url))

    # Mốc lề trái "chuẩn" của TOÀN VĂN BẢN = các giá trị x0 xuất hiện nhiều
    # lần. Dùng để phân biệt 1 khối THỰC SỰ là đầu câu hỏi mới với 1 mảnh
    # phân số (tử/mẫu) bị PDF tách thành block riêng nằm lệch hẳn giữa dòng
    # (ví dụ "7." đứng một mình do công thức y = 1/7 bị tách làm 2 block).
    common_margins = {x for x, cnt in margin_counter.items() if cnt >= 2}

    def near_margin(x0):
        if not common_margins:
            return True
        bucket = round(x0 / 4) * 4
        return any(abs(bucket - m) <= 12 for m in common_margins)

    # ---- Bước 2: tách câu hỏi XUYÊN SUỐT mọi trang (không reset theo trang) ----
    q_blocks = []
    current_q = None
    for (page_idx, text, x0, y0, y1) in all_kept_blocks:
        cleaned = _join_wrapped_lines(text)
        if _looks_like_header_or_footer(cleaned):
            continue

        lines = text.splitlines()
        if not lines:
            continue

        m = _QUESTION_START_RE.match(lines[0])
        starts_like_option = lines[0].strip().startswith(("A.", "B.", "C.", "D.", "a)", "b)", "c)", "d)"))
        is_new_question = bool(m) and near_margin(x0) and not starts_like_option

        if is_new_question:
            if current_q:
                q_blocks.append(current_q)
            current_q = {
                "number": m.group(1),
                "start_page": page_idx,
                "y0": y0,
                "question_text_lines": [m.group(2)] if m.group(2).strip() else lines[1:],
                "option_blocks": [],
            }
        elif current_q is not None:
            has_option_marker = any(_OPTION_START_RE.match(_clean_extracted_line(ln)) for ln in lines)
            if not current_q["option_blocks"] and not has_option_marker:
                current_q["question_text_lines"].extend(lines)
            else:
                current_q["option_blocks"].append(text)
        # else: nội dung trước câu hỏi đầu tiên của cả đề (tiêu đề, mục tiêu...)
        # -> bỏ qua, không thuộc câu hỏi nào cả

    if current_q:
        q_blocks.append(current_q)

    # ---- Bước 3: gán ảnh cho câu hỏi theo (trang, toạ độ Y) — hỗ trợ đúng cả
    # câu bị tràn qua trang sau, rồi phân loại MCQ/tự luận ----
    all_questions = []
    for q_idx, q in enumerate(q_blocks):
        if q_idx + 1 < len(q_blocks):
            next_page, next_y0 = q_blocks[q_idx + 1]["start_page"], q_blocks[q_idx + 1]["y0"]
        else:
            next_page, next_y0 = 10 ** 9, 0

        q_imgs = []
        for (img_page, iy0, iy1, url) in all_page_images:
            if img_page < q["start_page"]:
                continue
            if img_page == q["start_page"] and iy0 < q["y0"] - 15:
                continue
            if img_page > next_page:
                continue
            if img_page == next_page and iy0 >= next_y0:
                continue
            q_imgs.append(url)

        combined_options = "\n".join(q["option_blocks"])
        parsed_options = {}
        for letter, value in _parse_option_lines(combined_options):
            if letter not in parsed_options or (value and not parsed_options[letter]):
                parsed_options[letter] = value

        qtext = _join_wrapped_lines("\n".join(q["question_text_lines"]))
        if not qtext:
            continue

        if all(k in parsed_options and parsed_options[k] for k in "ABCD"):
            all_questions.append(_make_question_dict(q["number"], qtext, "mcq", parsed_options, q["start_page"], q_imgs))
        elif not parsed_options:
            # Câu tự luận: có số thứ tự, có ảnh/bảng biến thiên đính kèm (nếu
            # có) nhưng không có 4 đáp án A/B/C/D đi kèm.
            all_questions.append(_make_question_dict(q["number"], qtext, "essay", {}, q["start_page"], q_imgs))
        # else: chỉ tìm thấy 1-3/4 đáp án -> khả năng cao là lỗi đọc, bỏ qua

    return all_questions


def _parse_questions_from_text_generic(raw_text: str):
    if not raw_text:
        return []

    lines = raw_text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    starts = []
    for i, line in enumerate(lines):
        m = _QUESTION_START_RE.match(line)
        if m and not line.strip().startswith(("A.", "B.", "C.", "D.", "a)", "b)", "c)", "d)")):
            starts.append((i, m.group(1), m.group(2)))

    questions = []
    for idx, (start, number, first_line) in enumerate(starts):
        end = starts[idx + 1][0] if idx + 1 < len(starts) else len(lines)
        block = lines[start:end]
        qtext_lines = [first_line]
        option_source = []
        seen_option = False

        for line in block[1:]:
            if _OPTION_START_RE.match(_clean_extracted_line(line)):
                seen_option = True
            if seen_option:
                option_source.append(line)
            else:
                qtext_lines.append(line)

        parsed_options = {}
        for letter, value in _parse_option_lines("\n".join(option_source)):
            if letter not in parsed_options or (value and not parsed_options[letter]):
                parsed_options[letter] = value

        qtext = _join_wrapped_lines(" ".join(qtext_lines))
        if not qtext:
            continue

        if all(k in parsed_options and parsed_options[k] for k in "ABCD"):
            questions.append(_make_question_dict(number, qtext, "mcq", parsed_options, None, []))
        elif not parsed_options:
            questions.append(_make_question_dict(number, qtext, "essay", {}, None, []))

    return questions


def _detect_answer_key(raw_text: str):
    if not raw_text:
        return {}
    answers = {}
    patterns = [
        r"(?:Đáp\s*án|Dap\s*an|Answer)\s*(?:câu\s*)?(\d+(?:\.\d+)?)\s*[:.)-]?\s*([ABCD])\b",
        r"(?:Câu\s*)?(\d+(?:\.\d+)?)\s*[:.)-]\s*(?:Đáp\s*án|Dap\s*an|Answer)\s*[:.)-]?\s*([ABCD])\b",
    ]
    for pattern in patterns:
        for m in re.finditer(pattern, raw_text, flags=re.IGNORECASE):
            answers[m.group(1)] = m.group(2).upper()

    key_section = re.search(r"đáp\s*án|dap\s*an|answer\s*key", raw_text, flags=re.IGNORECASE)
    if key_section:
        tail = raw_text[key_section.end():]
        for m in re.finditer(r"\b(\d+(?:\.\d+)?)\s*[.\-:)]\s*([ABCD])\b", tail):
            answers.setdefault(m.group(1), m.group(2).upper())

    return answers


def _extract_questions_from_xlsx(file_path: str):
    """Đọc câu hỏi/đáp án trực tiếp từ ô Excel. Hỗ trợ 2 kiểu:
    1) Bảng cột: câu hỏi | A | B | C | D | (đáp án đúng) -> trắc nghiệm.
    2) Chỉ có câu hỏi, không có 4 đáp án ở 4 ô kế tiếp -> tự luận."""
    if not OPENPYXL_AVAILABLE:
        return []
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    except Exception:
        return []

    questions = []
    free_text_blocks = []
    auto_number = 0

    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            non_empty = [c for c in cells if c]
            if not non_empty:
                continue

            first_lower = non_empty[0].lower()
            if first_lower in {"câu hỏi", "cau hoi", "question", "stt", "số", "so", "câu", "no."}:
                continue
            if len(cells) >= 5 and all(
                (cells[i] or "").strip().lower() == letter
                for i, letter in enumerate(("a", "b", "c", "d"), start=1)
            ):
                continue

            if len(cells) >= 5 and cells[0] and cells[1] and cells[2] and cells[3] and cells[4]:
                auto_number += 1
                m = _QUESTION_START_RE.match(cells[0])
                if m:
                    q_number, q_text = m.group(1), m.group(2)
                else:
                    q_number, q_text = str(auto_number), cells[0]

                correct = None
                if len(cells) >= 6 and cells[5]:
                    letter_match = re.match(r"^\s*([ABCD])\b", cells[5], re.IGNORECASE)
                    if letter_match:
                        correct = letter_match.group(1).upper()

                questions.append({
                    "question_text": f"Câu {q_number}. {q_text}",
                    "question_type": "mcq",
                    "option_a": cells[1], "option_b": cells[2],
                    "option_c": cells[3], "option_d": cells[4],
                    "correct_answer": correct,
                    "correct_answer_text": None,
                    "source_page": ws.title,
                    "image_urls": "[]",
                })
            elif cells[0] and not any((cells[i] if len(cells) > i else "") for i in range(1, 5)):
                auto_number += 1
                m = _QUESTION_START_RE.match(cells[0])
                if m:
                    q_number, q_text = m.group(1), m.group(2)
                else:
                    q_number, q_text = str(auto_number), cells[0]
                questions.append(_make_question_dict(q_number, q_text, "essay", {}, ws.title, []))
            else:
                free_text_blocks.extend(non_empty)

    free_text_questions = _parse_questions_from_text_generic("\n".join(free_text_blocks)) if free_text_blocks else []
    return questions + free_text_questions


def extract_questions_universal(file_path: str, exam_id: int = None):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _extract_questions_from_pdf_layout(file_path, exam_id=exam_id) if PDF_AVAILABLE else []

    if ext in {".docx", ".doc", ".pptx", ".ppt"}:
        pdf_path = _convert_office_to_pdf(file_path)
        if pdf_path and PDF_AVAILABLE:
            doc_pdf_text = "\n".join(
                p.get_text("text", sort=True) for p in fitz.open(pdf_path)
            )
            if not _looks_like_broken_font_rendering(doc_pdf_text):
                questions = _extract_questions_from_pdf_layout(pdf_path, exam_id=exam_id)
                if questions:
                    return questions
        raw = _read_docx(file_path, exam_id=exam_id) if ext in (".docx", ".doc") else _read_pptx(file_path)
        return _parse_questions_from_text_generic(raw) if raw else []

    if ext in {".xlsx", ".xls"}:
        return _extract_questions_from_xlsx(file_path)

    if ext == ".txt":
        raw = _read_txt(file_path)
        return _parse_questions_from_text_generic(raw) if raw else []

    raw = extract_text_from_file(file_path, exam_id=exam_id)
    return _parse_questions_from_text_generic(raw) if raw else []


def parse_questions_pipeline(raw_text: str, file_path: str = None, exam_id: int = None):
    questions = []
    if file_path:
        questions = extract_questions_universal(file_path, exam_id=exam_id)

    if not questions and raw_text and raw_text.strip():
        questions = _parse_questions_from_text_generic(raw_text)

    if not questions:
        return []

    explicit_answers = _detect_answer_key(raw_text or "")
    for q in questions:
        if q.get("correct_answer"):
            continue
        number_match = re.search(r"(\d+(?:\.\d+)?)", q["question_text"])
        number = number_match.group(1) if number_match else None
        q["correct_answer"] = explicit_answers.get(number)

    return questions


def parse_questions_from_text_regex(raw_text: str):
    return _parse_questions_from_text_generic(raw_text)

# --------------------------------------------------------------------------
# Randomizer & Grading
# --------------------------------------------------------------------------

_GRADE_SUFFIX_RE = re.compile(r"(10|11|12)\s*$")


def _parse_subject_grade(subject_text):
    """Tách chuỗi giáo viên nhập theo cú pháp 'Môn Lớp' (VD: 'Toán 12',
    'Hóa 10', 'Lí 11') thành (tên môn, lớp). Không phụ thuộc danh sách môn cố
    định — bất kỳ tên môn nào có số lớp 10/11/12 ở cuối đều tách được, để
    dùng chung cho mọi môn học, không riêng Toán/Lý/Hóa/Văn."""
    if not subject_text:
        return (None, None)
    text = subject_text.strip()
    m = _GRADE_SUFFIX_RE.search(text)
    if m:
        grade = m.group(1)
        name = text[:m.start()].strip()
        return (name or text, grade)
    return (text, None)


def _group_items_by_subject_grade(items):
    """Gom danh sách (đề thi/tài liệu) có trường 'subject' dạng 'Môn Lớp'
    thành cấu trúc lồng: mỗi môn -> các lớp -> danh sách item. Item không có
    môn/lớp rõ ràng được gom vào nhóm 'Khác' để không bị mất."""
    groups = {}
    order = []
    for it in items:
        subj_raw = it["subject"] if "subject" in it.keys() else None
        name, grade = _parse_subject_grade(subj_raw)
        if not name:
            name, grade = "Khác", None
        if name not in groups:
            groups[name] = {"name": name, "grades": {}}
            order.append(name)
        grade_key = grade or "Khác"
        groups[name]["grades"].setdefault(grade_key, []).append(it)

    result = []
    for name in order:
        g = groups[name]
        grade_keys = sorted(g["grades"].keys(), key=lambda x: (x == "Khác", x))
        result.append({
            "name": name,
            "count": sum(len(v) for v in g["grades"].values()),
            "grades": [{"grade": gk, "items": g["grades"][gk]} for gk in grade_keys],
        })

    result.sort(key=lambda g: (g["name"] == "Khác", g["name"]))
    return result


def _group_items_by_grade_subject(items):
    """Gom theo LỚP trước rồi tới MÔN trong từng lớp (Lớp 10/11/12 -> Toán/Lí/
    Hóa/Văn...) — dùng riêng cho cây điều hướng ở sidebar, để học sinh chọn
    đúng khối lớp của mình trước rồi mới lọc theo môn."""
    groups = {}
    order = []
    for it in items:
        subj_raw = it["subject"] if "subject" in it.keys() else None
        name, grade = _parse_subject_grade(subj_raw)
        if not name:
            name, grade = "Khác", None
        grade_key = grade or "Khác"
        if grade_key not in groups:
            groups[grade_key] = {"grade": grade_key, "subjects": {}}
            order.append(grade_key)
        groups[grade_key]["subjects"].setdefault(name, []).append(it)

    grade_order = sorted(order, key=lambda g: (g == "Khác", g))
    result = []
    for gk in grade_order:
        g = groups[gk]
        subj_names = sorted(g["subjects"].keys(), key=lambda n: (n == "Khác", n))
        result.append({
            "grade": gk,
            "count": sum(len(v) for v in g["subjects"].values()),
            "subjects": [{"name": sn, "items": g["subjects"][sn]} for sn in subj_names],
        })
    return result


def seeded_random(exam_id: int, student_code: str) -> random.Random:
    key = f"{exam_id}:{student_code.strip().lower()}".encode("utf-8")
    seed = int(hashlib.sha256(key).hexdigest(), 16) % (2**32)
    return random.Random(seed)


def generate_test(exam_id: int, student_name: str, student_code: str, num_questions: int = None):
    db = get_db()
    all_questions = db.execute(
        "SELECT * FROM questions WHERE exam_id = ? ORDER BY order_index", (exam_id,)
    ).fetchall()
    if not all_questions:
        return None

    rng = seeded_random(exam_id, student_code)
    q_ids = [q["id"] for q in all_questions]

    if RANDOMIZE_EXAM:
        rng.shuffle(q_ids)

    if num_questions:
        q_ids = q_ids[:num_questions]

    option_maps = {}
    for q in all_questions:
        if q["id"] not in q_ids:
            continue
        if q["question_type"] == "essay":
            continue  # câu tự luận không có 4 phương án để xáo trộn
        letters = ["A", "B", "C", "D"]
        if RANDOMIZE_EXAM:
            shuffled = letters[:]
            rng.shuffle(shuffled)
        else:
            shuffled = letters[:]
        option_maps[str(q["id"])] = {display: original for display, original in zip(letters, shuffled)}

    cur = db.execute(
        "INSERT INTO tests (exam_id, student_name, student_code, question_order, option_maps, created_at) "
        "VALUES (?, ?, ?, ?, ?, ?)",
        (
            exam_id,
            student_name,
            student_code,
            json.dumps(q_ids),
            json.dumps(option_maps),
            datetime.now().isoformat(timespec="seconds"),
        ),
    )
    db.commit()
    return cur.lastrowid


def _normalize_answer_text(s: str) -> str:
    """Chuẩn hoá đáp án tự luận trước khi so khớp: bỏ khoảng trắng thừa, viết
    thường, gộp dấu trừ Unicode (−) về dấu gạch ngang thường (-), bỏ dấu chấm
    cuối câu — để 'trùng khớp' không đòi hỏi gõ chính xác từng ký tự/hoa
    thường như đáp án gốc."""
    if s is None:
        return ""
    s = s.strip().lower()
    s = s.replace("−", "-").replace("–", "-")
    s = re.sub(r"\s+", " ", s)
    s = s.rstrip(".")
    return s.strip()


def grade_submission(test_row, answers: dict):
    db = get_db()
    q_ids = json.loads(test_row["question_order"])
    option_maps = json.loads(test_row["option_maps"])

    questions = {
        q["id"]: q
        for q in db.execute(
            f"SELECT * FROM questions WHERE id IN ({','.join('?' * len(q_ids))})", q_ids
        ).fetchall()
    }

    results = []
    score = 0
    total_gradable = 0
    for qid in q_ids:
        q = questions[qid]

        if q["question_type"] == "essay":
            student_answer = (answers.get(str(qid)) or "").strip()
            has_key = bool(q["correct_answer_text"] and q["correct_answer_text"].strip())
            is_correct = None
            if has_key:
                total_gradable += 1
                is_correct = _normalize_answer_text(student_answer) == _normalize_answer_text(q["correct_answer_text"])
                if is_correct:
                    score += 1
            results.append({
                "question_text": q["question_text"],
                "question_type": "essay",
                "essay_answer": student_answer,
                "correct_answer_text": q["correct_answer_text"],
                "has_key": has_key,
                "options": {},
                "chosen": None,
                "correct_display": None,
                "is_correct": is_correct,
                "image_urls": q["image_urls"],
            })
            continue

        total_gradable += 1
        display_map = option_maps[str(qid)]
        chosen_display = answers.get(str(qid))
        chosen_original = display_map.get(chosen_display) if chosen_display else None
        is_correct = chosen_original is not None and chosen_original == q["correct_answer"]
        if is_correct:
            score += 1

        orig_texts = {
            "A": q["option_a"], "B": q["option_b"], "C": q["option_c"], "D": q["option_d"]
        }
        displayed_options = {display: orig_texts[orig] for display, orig in display_map.items()}
        correct_display = next((d for d, orig in display_map.items() if orig == q["correct_answer"]), None)

        results.append(
            {
                "question_text": q["question_text"],
                "question_type": "mcq",
                "options": displayed_options,
                "chosen": chosen_display,
                "correct_display": correct_display,
                "is_correct": is_correct,
                "image_urls": q["image_urls"],
            }
        )

    return score, total_gradable, results


# --------------------------------------------------------------------------
# Web Routes
# --------------------------------------------------------------------------

@app.route("/")
def home():
    db = get_db()
    exams = db.execute(
        "SELECT * FROM exams WHERE status = 'published' ORDER BY created_at DESC"
    ).fetchall()
    today = date.today().isoformat()
    today_exams = [e for e in exams if not e["available_date"] or e["available_date"] <= today]
    upcoming_exams = [e for e in exams if e["available_date"] and e["available_date"] > today]
    subject_groups = _group_items_by_subject_grade(today_exams)
    grade_groups = _group_items_by_grade_subject(today_exams)
    return render_template(
        "home.html", center_name=CENTER_NAME, today_exams=today_exams, upcoming_exams=upcoming_exams,
        subject_groups=subject_groups, grade_groups=grade_groups, active="kiemtra",
    )


@app.route("/exam/<int:exam_id>")
def exam_entry(exam_id):
    db = get_db()
    exam = db.execute("SELECT * FROM exams WHERE id = ? AND status='published'", (exam_id,)).fetchone()
    if not exam:
        abort(404)
    today = date.today().isoformat()
    if exam["available_date"] and exam["available_date"] > today:
        flash(f"Bài kiểm tra này chưa mở, sẽ mở vào ngày {exam['available_date']}.")
        return redirect(url_for("home"))
    return render_template("entry.html", center_name=CENTER_NAME, exam=exam, active="kiemtra")


@app.route("/exam/<int:exam_id>/start", methods=["POST"])
def exam_start(exam_id):
    db = get_db()
    exam = db.execute("SELECT * FROM exams WHERE id = ? AND status='published'", (exam_id,)).fetchone()
    if not exam:
        abort(404)
    today = date.today().isoformat()
    if exam["available_date"] and exam["available_date"] > today:
        flash(f"Bài kiểm tra này chưa mở, sẽ mở vào ngày {exam['available_date']}.")
        return redirect(url_for("home"))

    student_name = request.form.get("student_name", "").strip()
    student_code = request.form.get("student_code", "").strip()
    if not student_name or not student_code:
        flash("Vui lòng nhập đầy đủ họ tên và mã sinh viên.")
        return redirect(url_for("exam_entry", exam_id=exam_id))

    test_id = generate_test(exam_id, student_name, student_code)
    if test_id is None:
        flash("Đề này chưa có câu hỏi nào, vui lòng báo giáo viên.")
        return redirect(url_for("home"))

    return redirect(url_for("take_test", test_id=test_id))


@app.route("/test/<int:test_id>")
def take_test(test_id):
    db = get_db()
    test = db.execute("SELECT * FROM tests WHERE id = ?", (test_id,)).fetchone()
    if not test:
        abort(404)
    exam = db.execute("SELECT * FROM exams WHERE id = ?", (test["exam_id"],)).fetchone()

    q_ids = json.loads(test["question_order"])
    option_maps = json.loads(test["option_maps"])
    questions = {
        q["id"]: q
        for q in db.execute(
            f"SELECT * FROM questions WHERE id IN ({','.join('?' * len(q_ids))})", q_ids
        ).fetchall()
    }

    ordered_questions = []
    for idx, qid in enumerate(q_ids, start=1):
        q = questions[qid]
        if q["question_type"] == "essay":
            ordered_questions.append({
                "index": idx, "id": qid, "question_text": q["question_text"],
                "question_type": "essay", "options": [], "image_urls": q["image_urls"],
            })
            continue
        display_map = option_maps[str(qid)]
        orig_texts = {"A": q["option_a"], "B": q["option_b"], "C": q["option_c"], "D": q["option_d"]}
        displayed = [(letter, orig_texts[orig]) for letter, orig in sorted(display_map.items())]
        ordered_questions.append(
            {
                "index": idx, 
                "id": qid, 
                "question_text": q["question_text"], 
                "question_type": "mcq",
                "options": displayed,
                "image_urls": q["image_urls"]
            }
        )

    return render_template(
        "test.html", center_name=CENTER_NAME, exam=exam, test=test, questions=ordered_questions,
        active="kiemtra",
    )


@app.route("/test/<int:test_id>/submit", methods=["POST"])
def submit_test(test_id):
    db = get_db()
    test = db.execute("SELECT * FROM tests WHERE id = ?", (test_id,)).fetchone()
    if not test:
        abort(404)

    q_ids = json.loads(test["question_order"])
    answers = {}
    for qid in q_ids:
        val = request.form.get(f"q_{qid}")
        if val:
            answers[str(qid)] = val

    score, total, results = grade_submission(test, answers)

    db.execute(
        "INSERT INTO submissions (test_id, answers, score, total, submitted_at) VALUES (?, ?, ?, ?, ?)",
        (test_id, json.dumps(answers), score, total, datetime.now().isoformat(timespec="seconds")),
    )
    db.commit()
    submission_id = db.execute("SELECT last_insert_rowid() AS id").fetchone()["id"]

    return redirect(url_for("result", submission_id=submission_id))


@app.route("/result/<int:submission_id>")
def result(submission_id):
    db = get_db()
    submission = db.execute("SELECT * FROM submissions WHERE id = ?", (submission_id,)).fetchone()
    if not submission:
        abort(404)
    test = db.execute("SELECT * FROM tests WHERE id = ?", (submission["test_id"],)).fetchone()
    exam = db.execute("SELECT * FROM exams WHERE id = ?", (test["exam_id"],)).fetchone()

    answers = json.loads(submission["answers"])
    _, _, results = grade_submission(test, answers)

    duration_str = None
    try:
        started = datetime.fromisoformat(test["created_at"])
        finished = datetime.fromisoformat(submission["submitted_at"])
        total_seconds = max(0, int((finished - started).total_seconds()))
        m, s = divmod(total_seconds, 60)
        duration_str = f"{m} phút {s} giây" if m else f"{s} giây"
    except (ValueError, TypeError):
        duration_str = None

    return render_template(
        "result.html",
        center_name=CENTER_NAME,
        exam=exam,
        test=test,
        submission=submission,
        results=results,
        duration_str=duration_str,
        active="kiemtra",
    )


@app.route("/verify-teacher-password", methods=["POST"])
def verify_teacher_pass():
    password_input = request.form.get("teacher_password", "")
    next_url = request.form.get("next_url") or url_for("admin_home")
    
    if password_input == TEACHER_PASSWORD:
        return redirect(next_url)
    else:
        flash("❌ Mật khẩu giáo viên không chính xác. Vui lòng thử lại!")
        return redirect(request.referrer or url_for("home"))


@app.route("/admin")
def admin_home():
    db = get_db()
    exams = db.execute("SELECT * FROM exams ORDER BY created_at DESC").fetchall()
    exam_stats = {}
    for e in exams:
        n_questions = db.execute("SELECT COUNT(*) c FROM questions WHERE exam_id=?", (e["id"],)).fetchone()["c"]
        n_submissions = db.execute(
            "SELECT COUNT(*) c FROM submissions s JOIN tests t ON s.test_id=t.id WHERE t.exam_id=?", (e["id"],)
        ).fetchone()["c"]
        exam_stats[e["id"]] = {"n_questions": n_questions, "n_submissions": n_submissions}
    return render_template(
        "admin_home.html", center_name=CENTER_NAME, exams=exams, exam_stats=exam_stats,
        ocr_available=(PADDLE_AVAILABLE or OCR_AVAILABLE), pdf_available=PDF_AVAILABLE,
        active="admin", admin_tab="exams",
    )


@app.route("/admin/upload", methods=["POST"])
def admin_upload():
    file = request.files.get("exam_image")
    title = request.form.get("title", "").strip() or "Đề chưa đặt tên"
    subject = request.form.get("subject", "").strip()
    available_date = request.form.get("available_date", "").strip() or None
    time_limit_raw = request.form.get("time_limit_minutes", "").strip()
    time_limit_minutes = int(time_limit_raw) if time_limit_raw.isdigit() and int(time_limit_raw) > 0 else None

    if not file or file.filename == "":
        flash("Vui lòng chọn tệp đề thi (PDF, DOCX, PPTX, XLSX, TXT hoặc Ảnh).")
        return redirect(url_for("admin_home"))

    allowed_extensions = {
        ".pdf", ".txt", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls",
        ".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tif", ".tiff",
    }
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in allowed_extensions:
        flash("Định dạng chưa được hỗ trợ. Hãy dùng PDF, TXT, DOCX, PPTX, XLSX hoặc ảnh.")
        return redirect(url_for("admin_home"))

    filename = secure_filename(file.filename)
    stamp = datetime.now().strftime("%Y%m%d%H%M%S")
    saved_name = f"{stamp}_{filename}"
    saved_path = os.path.join(UPLOAD_DIR, saved_name)
    file.save(saved_path)

    db = get_db()
    cur = db.execute(
        "INSERT INTO exams (title, subject, source_image, raw_ocr_text, status, available_date, "
        "time_limit_minutes, created_at) VALUES (?, ?, ?, '', 'draft', ?, ?, ?)",
        (title, subject, saved_name, available_date, time_limit_minutes,
         datetime.now().isoformat(timespec="seconds")),
    )
    exam_id = cur.lastrowid

    raw_text = extract_text_from_file(saved_path, exam_id=exam_id)
    db.execute("UPDATE exams SET raw_ocr_text = ? WHERE id = ?", (raw_text, exam_id))

    parsed = parse_questions_pipeline(raw_text, saved_path, exam_id=exam_id) if raw_text else []

    for idx, q in enumerate(parsed):
        db.execute(
            "INSERT INTO questions (exam_id, order_index, question_text, question_type, option_a, "
            "option_b, option_c, option_d, correct_answer, correct_answer_text, image_urls) "
            "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
            (
                exam_id, idx, q["question_text"], q.get("question_type", "mcq"),
                q["option_a"], q["option_b"], q["option_c"], q["option_d"],
                q["correct_answer"], q.get("correct_answer_text"), q.get("image_urls", "[]"),
            ),
        )
    db.commit()

    n_essay = sum(1 for q in parsed if q.get("question_type") == "essay")
    if parsed:
        extra = f" (trong đó {n_essay} câu tự luận)" if n_essay else ""
        flash(f"Hệ thống đã tự động trích xuất thành công {len(parsed)} câu hỏi từ tệp {filename}!{extra}")
    else:
        flash("Không thể tự động nhận diện câu hỏi từ tệp này. Bạn có thể thêm câu hỏi thủ công ở trang bên dưới.")

    return redirect(url_for("admin_edit_exam", exam_id=exam_id))


@app.route("/admin/exam/<int:exam_id>/delete", methods=["POST"])
def admin_delete_exam(exam_id):
    db = get_db()
    exam = db.execute("SELECT * FROM exams WHERE id=?", (exam_id,)).fetchone()
    if exam:
        if exam["source_image"]:
            file_path = os.path.join(UPLOAD_DIR, exam["source_image"])
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except OSError:
                    pass
        db.execute("DELETE FROM exams WHERE id=?", (exam_id,))
        db.commit()
        flash(f"Đã xóa thành công đề thi '{exam['title']}'.")
    return redirect(url_for("admin_home"))


@app.route("/admin/exam/<int:exam_id>/edit")
def admin_edit_exam(exam_id):
    db = get_db()
    exam = db.execute("SELECT * FROM exams WHERE id=?", (exam_id,)).fetchone()
    if not exam:
        abort(404)
    questions = db.execute(
        "SELECT * FROM questions WHERE exam_id=? ORDER BY order_index", (exam_id,)
    ).fetchall()
    return render_template(
        "admin_edit.html", center_name=CENTER_NAME, exam=exam, questions=questions, active="admin", admin_tab="exams",
    )


@app.route("/admin/exam/<int:exam_id>/schedule", methods=["POST"])
def admin_update_schedule(exam_id):
    db = get_db()
    available_date = request.form.get("available_date", "").strip() or None
    db.execute("UPDATE exams SET available_date=? WHERE id=?", (available_date, exam_id))
    db.commit()
    return redirect(url_for("admin_edit_exam", exam_id=exam_id))


@app.route("/admin/exam/<int:exam_id>/time_limit", methods=["POST"])
def admin_update_time_limit(exam_id):
    db = get_db()
    raw = request.form.get("time_limit_minutes", "").strip()
    time_limit_minutes = int(raw) if raw.isdigit() and int(raw) > 0 else None
    db.execute("UPDATE exams SET time_limit_minutes=? WHERE id=?", (time_limit_minutes, exam_id))
    db.commit()
    return redirect(url_for("admin_edit_exam", exam_id=exam_id))


@app.route("/admin/exam/<int:exam_id>/add_question", methods=["POST"])
def admin_add_question(exam_id):
    db = get_db()
    q_type = request.form.get("question_type", "mcq")
    if q_type not in ("mcq", "essay"):
        q_type = "mcq"
    max_idx = db.execute(
        "SELECT COALESCE(MAX(order_index), -1) m FROM questions WHERE exam_id=?", (exam_id,)
    ).fetchone()["m"]
    default_correct = "A" if q_type == "mcq" else None
    db.execute(
        "INSERT INTO questions (exam_id, order_index, question_text, question_type, option_a, "
        "option_b, option_c, option_d, correct_answer, correct_answer_text, image_urls) "
        "VALUES (?, ?, '', ?, '', '', '', '', ?, NULL, '[]')",
        (exam_id, max_idx + 1, q_type, default_correct),
    )
    db.commit()
    return redirect(url_for("admin_edit_exam", exam_id=exam_id))


@app.route("/admin/question/<int:question_id>/toggle_type", methods=["POST"])
def admin_toggle_question_type(question_id):
    db = get_db()
    q = db.execute("SELECT * FROM questions WHERE id=?", (question_id,)).fetchone()
    if not q:
        abort(404)
    new_type = "essay" if q["question_type"] == "mcq" else "mcq"
    if new_type == "mcq":
        db.execute(
            "UPDATE questions SET question_type=?, correct_answer=COALESCE(correct_answer, 'A') WHERE id=?",
            (new_type, question_id),
        )
    else:
        db.execute(
            "UPDATE questions SET question_type=?, correct_answer=NULL WHERE id=?",
            (new_type, question_id),
        )
    db.commit()
    return redirect(url_for("admin_edit_exam", exam_id=q["exam_id"]))


@app.route("/admin/question/<int:question_id>/update", methods=["POST"])
def admin_update_question(question_id):
    db = get_db()
    q = db.execute("SELECT * FROM questions WHERE id=?", (question_id,)).fetchone()
    if not q:
        abort(404)
    correct = request.form.get("correct_answer", "").strip().upper() or None
    correct_text = request.form.get("correct_answer_text", "").strip() or None
    db.execute(
        "UPDATE questions SET question_text=?, option_a=?, option_b=?, option_c=?, option_d=?, "
        "correct_answer=?, correct_answer_text=? WHERE id=?",
        (
            request.form.get("question_text", "").strip(),
            request.form.get("option_a", "").strip(),
            request.form.get("option_b", "").strip(),
            request.form.get("option_c", "").strip(),
            request.form.get("option_d", "").strip(),
            correct,
            correct_text,
            question_id,
        ),
    )
    db.commit()
    return redirect(url_for("admin_edit_exam", exam_id=q["exam_id"]))


@app.route("/admin/exam/<int:exam_id>/save_all_questions", methods=["POST"])
def admin_save_all_questions(exam_id):
    db = get_db()
    question_ids = request.form.getlist("question_ids")

    for qid in question_ids:
        q_text = request.form.get(f"question_text_{qid}", "").strip()
        opt_a = request.form.get(f"option_a_{qid}", "").strip()
        opt_b = request.form.get(f"option_b_{qid}", "").strip()
        opt_c = request.form.get(f"option_c_{qid}", "").strip()
        opt_d = request.form.get(f"option_d_{qid}", "").strip()
        correct = request.form.get(f"correct_answer_{qid}", "").strip().upper() or None
        correct_text = request.form.get(f"correct_answer_text_{qid}", "").strip() or None

        db.execute(
            """
            UPDATE questions
            SET question_text=?, option_a=?, option_b=?, option_c=?, option_d=?,
                correct_answer=?, correct_answer_text=?
            WHERE id=? AND exam_id=?
            """,
            (q_text, opt_a, opt_b, opt_c, opt_d, correct, correct_text, qid, exam_id)
        )

    db.commit()
    flash(f"Đã lưu thành công toàn bộ {len(question_ids)} câu hỏi!")
    return redirect(url_for("admin_edit_exam", exam_id=exam_id))


@app.route("/admin/question/<int:question_id>/delete", methods=["POST"])
def admin_delete_question(question_id):
    db = get_db()
    q = db.execute("SELECT * FROM questions WHERE id=?", (question_id,)).fetchone()
    if not q:
        abort(404)
    db.execute("DELETE FROM questions WHERE id=?", (question_id,))
    db.commit()
    return redirect(url_for("admin_edit_exam", exam_id=q["exam_id"]))

@app.route("/admin/question/<int:question_id>/delete_image", methods=["POST"])
def admin_delete_question_image(question_id):
    db = get_db()
    q = db.execute("SELECT * FROM questions WHERE id=?", (question_id,)).fetchone()
    if not q:
        abort(404)
        
    image_url_to_delete = request.form.get("image_url", "").strip()
    if image_url_to_delete and q["image_urls"]:
        try:
            current_images = json.loads(q["image_urls"])
            # Lọc bỏ ảnh mà user chọn xóa
            updated_images = [img for img in current_images if img != image_url_to_delete]
            
            # Cập nhật lại danh sách ảnh vào Database
            db.execute("UPDATE questions SET image_urls=? WHERE id=?", (json.dumps(updated_images), question_id))
            db.commit()
            
            # (Tùy chọn) Xóa file thực tế trên ổ đĩa nếu là file tĩnh cục bộ
            clean_rel_path = image_url_to_delete.lstrip("/")
            file_disk_path = os.path.join(BASE_DIR, clean_rel_path)
            if os.path.exists(file_disk_path):
                try:
                    os.remove(file_disk_path)
                except OSError:
                    pass
            flash("Đã xóa hình ảnh thành công!")
        except Exception as e:
            print(f"[Delete Image Error] {e}")
            
    return redirect(url_for("admin_edit_exam", exam_id=q["exam_id"]))


@app.route("/admin/exam/<int:exam_id>/publish", methods=["POST"])
def admin_publish_exam(exam_id):
    db = get_db()
    n = db.execute("SELECT COUNT(*) c FROM questions WHERE exam_id=?", (exam_id,)).fetchone()["c"]
    if n == 0:
        flash("Đề chưa có câu hỏi nào, không thể đăng bài.")
        return redirect(url_for("admin_edit_exam", exam_id=exam_id))
    
    db.execute("UPDATE exams SET status='published' WHERE id=?", (exam_id,))
    db.commit()
    flash("Đã xuất bản bài kiểm tra thành công cho học sinh!")
    return redirect(url_for("admin_home"))


@app.route("/admin/exam/<int:exam_id>/unpublish", methods=["POST"])
def admin_unpublish_exam(exam_id):
    db = get_db()
    db.execute("UPDATE exams SET status='draft' WHERE id=?", (exam_id,))
    db.commit()
    return redirect(url_for("admin_home"))


@app.route("/admin/exam/<int:exam_id>/submissions")
def admin_exam_submissions(exam_id):
    db = get_db()
    exam = db.execute("SELECT * FROM exams WHERE id=?", (exam_id,)).fetchone()
    if not exam:
        abort(404)
    rows = db.execute(
        """
        SELECT t.student_name, t.student_code, s.score, s.total, s.submitted_at
        FROM submissions s JOIN tests t ON s.test_id = t.id
        WHERE t.exam_id = ?
        ORDER BY s.submitted_at DESC
        """,
        (exam_id,),
    ).fetchall()
    return render_template(
        "admin_submissions.html", center_name=CENTER_NAME, exam=exam, rows=rows, active="admin", admin_tab="exams",
    )


@app.route("/admin/documents")
def admin_documents():
    db = get_db()
    docs = db.execute("SELECT * FROM documents ORDER BY uploaded_at DESC").fetchall()
    return render_template(
        "admin_documents.html", center_name=CENTER_NAME, docs=docs, active="admin", admin_tab="documents",
    )


@app.route("/documents")
def public_documents():
    db = get_db()
    docs = db.execute("SELECT * FROM documents ORDER BY uploaded_at DESC").fetchall()
    subject_groups = _group_items_by_subject_grade(docs)
    grade_groups = _group_items_by_grade_subject(docs)
    return render_template(
        "documents.html", center_name=CENTER_NAME, docs=docs, subject_groups=subject_groups,
        grade_groups=grade_groups, active="tailieu",
    )


@app.route("/info")
def info_page():
    db = get_db()
    history_rows = db.execute(
        """
        SELECT 
            s.id AS submission_id,
            t.student_name,
            t.student_code,
            e.title AS exam_title,
            e.subject,
            s.score,
            s.total,
            t.created_at AS started_at,
            s.submitted_at
        FROM submissions s
        JOIN tests t ON s.test_id = t.id
        JOIN exams e ON t.exam_id = e.id
        ORDER BY s.submitted_at DESC
        """
    ).fetchall()

    processed_history = []
    for row in history_rows:
        duration_str = "-"
        try:
            started = datetime.fromisoformat(row["started_at"])
            finished = datetime.fromisoformat(row["submitted_at"])
            total_seconds = max(0, int((finished - started).total_seconds()))
            m, s = divmod(total_seconds, 60)
            duration_str = f"{m} phút {s} giây" if m else f"{s} giây"
        except (ValueError, TypeError):
            pass

        processed_history.append({
            "submission_id": row["submission_id"],
            "student_name": row["student_name"],
            "student_code": row["student_code"],
            "exam_title": row["exam_title"],
            "subject": row["subject"],
            "score": row["score"],
            "total": row["total"],
            "submitted_at": row["submitted_at"],
            "duration": duration_str
        })

    return render_template(
        "info.html", 
        center_name=CENTER_NAME, 
        history=processed_history, 
        active="thongtin"
    )


@app.route("/info/submission/<int:submission_id>/delete", methods=["POST"])
def delete_submission(submission_id):
    db = get_db()
    sub = db.execute("SELECT * FROM submissions WHERE id = ?", (submission_id,)).fetchone()
    if sub:
        db.execute("DELETE FROM submissions WHERE id = ?", (submission_id,))
        db.execute("DELETE FROM tests WHERE id = ?", (sub["test_id"],))
        db.commit()
        flash("Đã xóa lượt làm bài của học sinh thành công.")
    return redirect(url_for("info_page"))


@app.route("/admin/documents/upload", methods=["POST"])
def admin_documents_upload():
    file = request.files.get("document_file")
    if not file or file.filename == "":
        flash("Vui lòng chọn tệp tài liệu.")
        return redirect(url_for("admin_documents"))

    subject = request.form.get("subject", "").strip() or None

    filename = secure_filename(file.filename)
    stamp = datetime.now().strftime("%Y%m%d%H%M%S")
    saved_name = f"{stamp}_{filename}"
    file.save(os.path.join(DOCS_DIR, saved_name))

    db = get_db()
    db.execute(
        "INSERT INTO documents (original_filename, stored_filename, subject, uploaded_at) VALUES (?, ?, ?, ?)",
        (filename, saved_name, subject, datetime.now().isoformat(timespec="seconds")),
    )
    db.commit()
    return redirect(url_for("admin_documents"))


@app.route("/admin/documents/<int:doc_id>/delete", methods=["POST"])
def admin_documents_delete(doc_id):
    db = get_db()
    doc = db.execute("SELECT * FROM documents WHERE id=?", (doc_id,)).fetchone()
    if doc:
        try:
            os.remove(os.path.join(DOCS_DIR, doc["stored_filename"]))
        except OSError:
            pass
        db.execute("DELETE FROM documents WHERE id=?", (doc_id,))
        db.commit()
    return redirect(url_for("admin_documents"))


@app.route("/documents/file/<path:filename>")
def serve_document(filename):
    return send_from_directory(DOCS_DIR, filename, as_attachment=True)


if __name__ == "__main__":
    init_db()
    if PDF_AVAILABLE:
        print("-> Document Parser (PyMuPDF) đã sẵn sàng.")
    if PADDLE_AVAILABLE:
        print("-> PaddleOCR đã sẵn sàng.")
    elif OCR_AVAILABLE:
        print("-> Tesseract OCR (dự phòng) đã sẵn sàng.")
    app.run(debug=True, use_reloader=False, port=5000)
